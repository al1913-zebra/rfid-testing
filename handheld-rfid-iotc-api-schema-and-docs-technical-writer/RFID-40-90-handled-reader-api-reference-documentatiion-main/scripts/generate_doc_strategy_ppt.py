from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("docs/handled_reader_usability_meeting.pptx")

BG = RGBColor(245, 247, 250)
NAVY = RGBColor(16, 44, 84)
TEAL = RGBColor(21, 125, 140)
GOLD = RGBColor(214, 160, 55)
TEXT = RGBColor(40, 46, 52)
MUTED = RGBColor(98, 108, 122)
WHITE = RGBColor(255, 255, 255)


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.34), Inches(0.65))
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.9), 0, Inches(1.44), Inches(0.65))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.95), Inches(11.2), Inches(0.8))
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = NAVY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.75), Inches(11), Inches(0.5))
        sub_frame = sub_box.text_frame
        sub_frame.word_wrap = True
        p = sub_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(12)
        r.font.color.rgb = MUTED


def add_bullets(slide, items, left=0.9, top=2.1, width=11.1, height=4.6, level0_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    first = True
    for item in items:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item

        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        paragraph.level = level
        paragraph.space_after = Pt(8)
        paragraph.line_spacing = 1.15
        run = paragraph.add_run()
        run.text = text
        run.font.name = "Aptos"
        run.font.size = Pt(level0_size if level == 0 else 15)
        run.font.color.rgb = TEXT if level == 0 else MUTED
        if level == 0:
            run.font.bold = True


def add_three_cards(slide, cards):
    positions = [0.72, 4.42, 8.12]
    for index, (title, lines, color) in enumerate(cards):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(positions[index]), Inches(2.05), Inches(3.05), Inches(4.25))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(positions[index]), Inches(2.05), Inches(3.05), Inches(0.55))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()

        head_tf = head.text_frame
        p = head_tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.name = "Aptos Display"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = WHITE

        add_bullets(slide, lines, left=positions[index] + 0.18, top=2.78, width=2.68, height=3.25, level0_size=16)


def add_two_column_table_like(slide, left_title, left_lines, right_title, right_lines):
    for index, (title, lines, x, color) in enumerate(
        [
            (left_title, left_lines, 0.7, NAVY),
            (right_title, right_lines, 6.75, TEAL),
        ]
    ):
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(5.55), Inches(4.4))
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = color
        panel.line.width = Pt(2)

        title_box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(2.18), Inches(5), Inches(0.45))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Aptos Display"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = color

        add_bullets(slide, lines, left=x + 0.25, top=2.75, width=5.0, height=3.2, level0_size=15)


def add_quote(slide, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.2), Inches(11.4), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(2)

    text_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.65), Inches(10.7), Inches(1.4))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos Display"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = NAVY


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(
        slide,
        "Handled Reader Documentation: Project Understanding",
        "A 3-slide meeting view of the documentation model and usability improvements",
    )
    add_three_cards(
        slide,
        [
            (
                "Event Commands",
                [
                    "Reader publishes payloads.",
                    "Docs focus on trigger, payload, and examples.",
                ],
                NAVY,
            ),
            (
                "Set / Config",
                [
                    "User changes parameters by requirement.",
                    "Docs focus on supported values, prerequisites, and rules.",
                ],
                TEAL,
            ),
            (
                "Get Commands",
                [
                    "Minimal input from user.",
                    "Docs focus on returned response and interpretation.",
                ],
                GOLD,
            ),
        ],
    )
    add_bullets(
        slide,
        [
            "Because command behavior differs, we avoided one generic template for every operation.",
            "`set_config` became the reference example for the complex configuration-style template.",
        ],
        top=6.45,
        height=0.65,
        level0_size=14,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Usability Improvements Implemented in the Docs Portal")
    add_two_column_table_like(
        slide,
        "Navigation & Discovery",
        [
            "Two search experiences: sidebar search for TOC filtering and topbar search for page-wide discovery.",
            "Categorized TOC groups operations by tag so users find the right section faster.",
            "Per-operation PDF download makes content easy to export and analyze with AI tools.",
        ],
        "Content Readability",
        [
            "JSON examples are presented in a structured interactive format with expand / collapse behavior.",
            "Schema is shown in table format with clear Field / Type / Description columns.",
            "Data types are color-coded so users can quickly distinguish enum, object, array, number, boolean, and string values.",
        ],
    )
    add_bullets(
        slide,
        [
            "This shifts the documentation from static reference to a usable exploration tool.",
        ],
        top=6.65,
        height=0.4,
        level0_size=15,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Why This Matters for the Product")
    add_three_cards(
        slide,
        [
            (
                "For Users",
                [
                    "Faster onboarding.",
                    "Less confusion while reading commands and schemas.",
                ],
                NAVY,
            ),
            (
                "For Support & Delivery",
                [
                    "Fewer clarifications on payload structure and supported values.",
                    "Better shareability through PDF export.",
                ],
                TEAL,
            ),
            (
                "For the Meeting",
                [
                    "Shows that the project is not only documented, but designed for usability.",
                    "Highlights both content quality and portal experience in one story.",
                ],
                GOLD,
            ),
        ],
    )
    add_quote(slide, "The goal is not only to document APIs, but to make them easier to understand, search, export, and use.")

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    return prs.slides


if __name__ == "__main__":
    build_presentation()
    print(f"Presentation written to {OUTPUT_PATH.resolve()}")